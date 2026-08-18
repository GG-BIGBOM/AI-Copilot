"""造测试用的样例文件（docx / pptx / pdf）。

**故意用真库真格式，不用 mock。** 上传解析这条链路上，坑几乎全在
「真实文件的结构和你想的不一样」——Word 的表格不在段落流里、
PPT 的版式没有标题占位符、扫描件 PDF 一个字也提不出来。
把这些造成假对象，测试就只在验证我自己的想象。
"""

from __future__ import annotations

from pathlib import Path


def make_docx(path: Path) -> Path:
    """一份带标题层级、正文和表格的 Word 文档。

    表格**夹在两段正文中间**——这是关键：python-docx 的 `paragraphs` 和
    `tables` 是两个独立列表，按它们拼出来的结果会把表格甩到文末。
    """
    import docx

    d = docx.Document()
    d.add_heading("一、电子面单设置", level=1)
    d.add_paragraph("先在设置里绑定物流账号。")
    d.add_heading("1、模板选择", level=2)
    d.add_paragraph("选择京东标准模板。")

    table = d.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "含义"
    table.cell(1, 0).text = "运单号"
    table.cell(1, 1).text = "物流公司生成的单号"

    d.add_paragraph("保存后即可打印。")
    d.save(str(path))
    return path


def make_pptx(path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # 标题 + 内容

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "面单打印流程"
    slide.placeholders[1].text = "第一步：绑定物流账号"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "常见问题"
    slide2.placeholders[1].text = "打印偏移怎么调"

    prs.save(str(path))
    return path


def make_pdf(path: Path, text: str = "Hello PDF world") -> Path:
    """手搓一个最小的、能提取出文字的 PDF。

    不用 reportlab 之类的生成库：那是给项目多加一个只有测试用得上的依赖，
    而 PDF 的结构简单到手写更清楚——也顺便说明了「PDF 里没有标题样式」
    这件事（见 parsers.parse_pdf 的取舍）。
    """
    body = f"BT /F1 12 Tf 20 100 Td ({text}) Tj ET".encode("latin-1")
    objs = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 200 200]/Contents 4 0 R"
        b"/Resources<</Font<</F1 5 0 R>>>>>>",
        b"<</Length %d>>stream\n" % len(body) + body + b"\nendstream",
        b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for i, obj in enumerate(objs, 1):
        offsets.append(len(out))
        out += b"%d 0 obj" % i + obj + b"endobj\n"

    xref = len(out)
    out += b"xref\n0 %d\n" % (len(objs) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF\n" % (len(objs) + 1, xref)

    path.write_bytes(bytes(out))
    return path


def make_scanned_pdf(path: Path) -> Path:
    """一页空白的 PDF，模拟扫描件——`extract_text()` 返回空串。"""
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with path.open("wb") as f:
        writer.write(f)
    return path


def make_png(path: Path, size=(120, 60), color=(240, 240, 240)) -> Path:
    """一张纯色 PNG。内容不重要——图片解析这条路上，"看懂什么"由视觉模型负责，
    测试要验的是它前后的管道：格式转换、失败处理、标题补齐。"""
    from PIL import Image

    Image.new("RGB", size, color).save(path, format="PNG")
    return path
