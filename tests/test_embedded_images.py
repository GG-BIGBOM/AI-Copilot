"""上传文档里的嵌图：DOCX / PPTX / XLSX / PDF（M17）。

这一组守四件事：

    1. **图跟着它说明的那一段走**，不是统一堆到文末。ERP 手册里图就是步骤
       本身——「第二步长什么样」答不了的话，解出来也没用。
    2. **归属正确**：第几页 / 第几张 slide / 哪个工作表。归属错了不会报错，
       只会让答案配上另一页的截图，而用户照着点会点不到。
    3. **私有图落在私有目录**（`data/private-images/`，nginx 不发它）。
       落进公共目录的表现是"图能显示"——**正是问题所在**：谁猜中文件名
       谁就能取别人 Word 里的截图。
    4. **闸门**：张数上限、单张大小上下限。一份 PPT 能塞几百张图，
       而 worker 的 MemoryMax 是 400M。

⚠️ fixture 全部在测试里现造（python-docx / python-pptx / openpyxl / 手写 PDF），
不进仓库：二进制文件进 Git 之后没人知道它里面到底有什么，而这几条断言
恰恰是关于「里面有什么」的。
"""

from __future__ import annotations

import io
import uuid
import zlib
from pathlib import Path

import pytest

from copilot import assets
from copilot.config import get_settings
from copilot.ingest.parsers import ParseError, parse_docx, parse_pdf, parse_pptx, parse_xlsx

# 1×1 PNG。真字节：解析器会把它原样交给 sink，sink 要能算 sha256、能落盘
PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000100fdff03fd0000000049454e44ae426082"
)


def _padded(seed: bytes = PNG, size: int = 8 * 1024) -> bytes:
    """撑到闸门下限以上的图。PNG 后面补零仍是合法 PNG（尾部数据被忽略）。

    ⚠️ 每张图的内容必须**不同**：存储是按内容寻址的，两张一模一样的图
    只会落一个文件、一行资产——那会让「三张图」这种断言默默地变成一张。
    """
    return seed + uuid.uuid4().bytes + b"\0" * (size - len(seed) - 16)


class Sink:
    """记账用的假 sink。只关心解析器**怎么调它**，不落盘。"""

    def __init__(self, limit: int = 99) -> None:
        self.calls: list[dict] = []
        self.limit = limit

    def save(self, data: bytes, suffix: str, **meta) -> str | None:
        if len(self.calls) >= self.limit:
            return None
        self.calls.append({"size": len(data), "suffix": suffix, **meta})
        return f"asset://fake/{len(self.calls)}{suffix}"


# ─────────────── DOCX ───────────────


def _docx(path: Path) -> Path:
    import docx

    d = docx.Document()
    d.add_heading("电子面单设置", level=1)
    d.add_paragraph("第一步：在【设置-物流】里绑定物流账号。")
    d.add_picture(io.BytesIO(_padded()))
    d.add_paragraph("第二步：回到订单列表，勾选后点打印。")
    d.add_picture(io.BytesIO(_padded()))
    d.save(str(path))
    return path


def test_docx_images_land_next_to_their_paragraph(tmp_path):
    """⭐ 图在它所属那一段的**后面**，不是文末。

    统一堆到文末的话，切分之后所有图都落在最后一个块里——「第二步长什么样」
    这个问题就再也答不了了。
    """
    sink = Sink()

    md = parse_docx(_docx(tmp_path / "a.docx"), sink=sink).markdown

    lines = [ln for ln in md.splitlines() if ln.strip()]
    first_img = next(i for i, ln in enumerate(lines) if ln.startswith("!["))
    assert "第一步" in lines[first_img - 1], f"图没跟着它那一段：\n{md}"
    assert md.count("![](asset://") == 2
    # 最后一段后面还有一张，说明两张图各自跟着各自的段落
    assert lines[-1].startswith("![]"), md


def test_docx_without_a_sink_is_unchanged(tmp_path):
    """⚠️ 不给 sink 就**完全照旧**：只出文字，一行不多。

    语雀那条入库路径、评测、CLI 都不带 sink，它们的行为一个字都不该变。
    """
    md = parse_docx(_docx(tmp_path / "b.docx")).markdown

    assert "![" not in md
    assert "第一步" in md and "第二步" in md


# ─────────────── PPTX ───────────────


def _pptx(path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for i, text in enumerate(("退货入库", "打印面单"), 1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = text
        # 第 2 页放两张图，第 1 页一张——用来验「绑到自己那一页」
        for _ in range(i):
            slide.shapes.add_picture(
                io.BytesIO(_padded()), Inches(1), Inches(2), Inches(1), Inches(1)
            )
    prs.save(str(path))
    return path


def test_pptx_images_are_bound_to_their_slide(tmp_path):
    """⭐ `slide_number` 必须是**它自己那一页**。绑错页 = 答案配错截图。"""
    sink = Sink()

    md = parse_pptx(_pptx(tmp_path / "a.pptx"), sink=sink).markdown

    assert [c["slide_number"] for c in sink.calls] == [1, 2, 2]
    # 正文里也要落在对应那一节下面
    first, second = md.split("## 第 2 页")
    assert first.count("![]") == 1
    assert second.count("![]") == 2


# ─────────────── XLSX ───────────────


def _xlsx(path: Path) -> Path:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    wb = Workbook()
    ws = wb.active
    ws.title = "对账规则"
    ws.append(["字段", "含义"])
    ws.append(["单号", "平台订单号"])

    img_path = path.parent / "cell.png"
    img_path.write_bytes(_padded())
    ws.add_image(XLImage(str(img_path)), "B7")

    wb.create_sheet("空表")
    wb.save(str(path))
    return path


def test_xlsx_becomes_a_table_and_keeps_its_sheet(tmp_path):
    """Excel：一个工作表一节，表格转成 Markdown 表格，图记住是哪个表的。"""
    sink = Sink()

    md = parse_xlsx(_xlsx(tmp_path / "a.xlsx"), sink=sink).markdown

    assert "## 对账规则" in md
    assert "| 单号 | 平台订单号 |" in md
    assert sink.calls and sink.calls[0]["sheet_name"] == "对账规则"
    assert sink.calls[0]["anchor"], "没记下它钉在哪个单元格"
    # 空工作表不该留下一个只有标题的空节
    assert "## 空表" not in md


def test_an_empty_workbook_is_an_error_not_an_empty_document(tmp_path):
    """⚠️ 空内容要报错，不能静默入库——用户会看到「已完成」却永远搜不到。"""
    from openpyxl import Workbook

    path = tmp_path / "empty.xlsx"
    Workbook().save(str(path))

    with pytest.raises(ParseError):
        parse_xlsx(path)


# ─────────────── PDF ───────────────


def _pdf(path: Path, text: str = "Hello ERP") -> Path:
    """手写一个「有文字 + 有嵌图」的最小 PDF。

    不用第三方写库：为了造一个测试文件引一个新依赖不划算，而 PDF 的
    这一小段结构十年没变过。
    """
    raw = bytes([255, 0, 0, 0, 255, 0, 0, 0, 255, 255, 255, 0]) * 700  # 撑过下限
    img = zlib.compress(raw)
    content = (
        f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET\nq 100 0 0 100 72 600 cm /Im1 Do Q\n"
    ).encode()

    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> /XObject << /Im1 6 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Type /XObject /Subtype /Image /Width 70 /Height 10 /ColorSpace /DeviceRGB "
        b"/BitsPerComponent 8 /Filter /FlateDecode /Length " + str(len(img)).encode() + b" >>\n"
        b"stream\n" + img + b"\nendstream",
    ]

    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, 1):
        offsets.append(out.tell())
        out.write(f"{i} 0 obj\n".encode() + body + b"\nendobj\n")
    xref = out.tell()
    out.write(f"xref\n0 {len(objs) + 1}\n".encode())
    out.write(b"0000000000 65535 f \n")
    for off in offsets:
        out.write(f"{off:010d} 00000 n \n".encode())
    out.write(
        f"trailer\n<< /Size {len(objs) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    path.write_bytes(out.getvalue())
    return path


def test_pdf_images_carry_their_page_number(tmp_path):
    sink = Sink()

    md = parse_pdf(_pdf(tmp_path / "a.pdf"), sink=sink).markdown

    assert "Hello ERP" in md
    assert sink.calls and sink.calls[0]["page_number"] == 1
    assert "![](asset://" in md


# ─────────────── 落盘：私有图必须在私有目录 ───────────────


def test_private_images_never_land_in_the_public_directory(tmp_path):
    """⭐⭐ 这一条是 M17 的安全落点。

    `data/images/` 是 nginx 直接 alias 出去的静态目录，谁猜中文件名谁就能取。
    别人 Word 里的截图落进去，表现是「图能正常显示」——**正是问题所在**，
    没有任何症状能提示你内容泄漏了。
    """
    s = get_settings()
    sink = assets.UploadImageSink(private=True)

    ref = sink.save(_padded(), ".png")

    assert ref and ref.startswith("asset://"), "私有图不能写成 /images/ 那种可直发的地址"
    rel = assets.storage_path_of(ref)
    assert (s.private_image_dir / rel).is_file(), "文件没落在私有目录"
    assert not (s.image_dir / rel).exists(), "私有图落进了 nginx 直发的公共目录"


def test_the_same_image_is_stored_once(tmp_path):
    """按内容寻址：同一张图重复出现只存一份，重传同一份文件也不翻倍。"""
    sink = assets.UploadImageSink(private=True)
    data = _padded()

    first, second = sink.save(data, ".png"), sink.save(data, ".png")

    assert first == second
    assert len(sink.saved) == 1


def test_the_gates_hold(tmp_path):
    """张数上限、太小、太大——三道闸门，被挡掉的返回 None 并记一笔。

    ⚠️ 返回 None 时解析器**不写这张图**：写一个取不到的地址比没有图更糟，
    页面上是一张裂图，而用户不知道为什么。
    """
    sink = assets.UploadImageSink(private=True, max_images=2)

    assert sink.save(_padded(), ".png")
    assert sink.save(_padded(), ".png")
    assert sink.save(_padded(), ".png") is None, "超出张数上限还在收"
    assert sink.save(b"tiny", ".png") is None, "图标那么小的图不该收"
    assert sink.save(b"x" * (get_settings().image_max_bytes + 1), ".png") is None
    assert sink.skipped == 3


def test_a_saved_image_is_addressable_by_the_api_path():
    """落盘之后，`storage_path_of` 要能把地址还原回磁盘路径——
    还原不回来的话，`sync_document_assets` 建不出资产行，
    而没有资产行的私有图会在检索层被丢掉（fail closed）。"""
    sink = assets.UploadImageSink(private=True)

    ref = sink.save(_padded(), ".jpg")

    rel = assets.storage_path_of(ref)
    assert rel and rel in sink.saved
    assert assets.mime_of(rel) == "image/jpeg"


# ─────────────── 端到端：上传 → 解析 → 只有本人取得到 → 删干净 ───────────────


@pytest.fixture
async def uploaded_docx(maker, logged_in, flagship_id, tmp_path):
    """走真实的 worker 路径解析一份带图的 docx。返回文档 id。"""
    import uuid as _uuid

    from chat_helpers import FakeEmbedder
    from sqlalchemy import delete

    from copilot.db.models import Chunk, Document, ImageAsset, Job
    from copilot.jobs import queue
    from copilot.jobs.worker import run_once

    s = get_settings()
    rel = f"{logged_in}/{_uuid.uuid4().hex}.docx"
    path = s.upload_path(rel)
    path.parent.mkdir(parents=True, exist_ok=True)
    _docx(path)

    async with maker() as session:
        doc = Document(
            owner_id=logged_in,
            knowledge_space_id=flagship_id,
            source_type="upload",
            title=f"带图手册-{_uuid.uuid4().hex[:8]}",
            original_filename="带图手册.docx",
            stored_path=rel,
            size_bytes=path.stat().st_size,
            content_hash=_uuid.uuid4().hex,
            status="pending",
        )
        session.add(doc)
        await session.flush()
        await queue.enqueue(session, queue.PARSE_UPLOAD, queue.document_payload(doc.id))
        await session.commit()
        doc_id = doc.id

    assert await run_once(FakeEmbedder(), maker) == "done"

    yield doc_id

    async with maker() as session:
        await session.execute(
            delete(Job).where(Job.payload["document_id"].astext == str(doc_id))
        )
        await session.execute(delete(ImageAsset).where(ImageAsset.document_id == doc_id))
        await session.execute(delete(Chunk).where(Chunk.document_id == doc_id))
        await session.execute(delete(Document).where(Document.id == doc_id))
        await session.commit()


async def test_an_uploaded_docx_ends_up_with_private_assets(maker, uploaded_docx, logged_in):
    """⭐ 真实链路：解析出的图有资产行，owner 跟着文档，文件在私有目录。"""
    from sqlalchemy import select

    from copilot.db.models import Chunk, Document, ImageAsset

    async with maker() as s:
        doc = await s.get(Document, uploaded_docx)
        assert doc.status == "done", doc.error
        rows = list(
            (
                await s.execute(select(ImageAsset).where(ImageAsset.document_id == uploaded_docx))
            ).scalars()
        )
        chunks = list(
            (await s.execute(select(Chunk).where(Chunk.document_id == uploaded_docx))).scalars()
        )

    assert len(rows) == 2, "两张图应该各有一行资产"
    for row in rows:
        # ⚠️ 隔离红线：资产的 owner 必须跟着文档，`/api/images/` 只看这一列
        assert row.owner_id == logged_in
        assert row.knowledge_space_id == doc.knowledge_space_id
        assert (get_settings().private_image_dir / row.storage_path).is_file()
        assert not (get_settings().image_dir / row.storage_path).exists()
        assert row.sha256 and row.file_size

    # 块上的对照表仍然是事实来源（双读兼容，M14-B）
    assert any(c.images for c in chunks), "块上没有配图对照表"


async def test_only_the_owner_can_fetch_an_embedded_image(
    api_client, maker, uploaded_docx, logged_in
):
    """⭐⭐ A 的图：A 取得到，别人和匿名一律 404。"""
    from sqlalchemy import delete, select

    from copilot.auth.security import create_access_token
    from copilot.db.models import ImageAsset, User

    async with maker() as s:
        asset_id = (
            await s.execute(
                select(ImageAsset.id).where(ImageAsset.document_id == uploaded_docx).limit(1)
            )
        ).scalar_one()

    mine = await api_client.get(f"/api/images/{asset_id}")
    assert mine.status_code == 200, mine.text
    assert mine.content[:8] == PNG[:8], "取回来的不是那张图"

    email = f"mallory-{uuid.uuid4().hex[:8]}@test.local"
    async with maker() as s:
        other = User(email=email, password_hash="x", is_active=True)
        s.add(other)
        await s.commit()
        token = create_access_token(other.id)

    try:
        api_client.cookies.clear()
        theirs = await api_client.get(
            f"/api/images/{asset_id}", headers={"Authorization": f"Bearer {token}"}
        )
        assert theirs.status_code == 404, theirs.text
        assert (await api_client.get(f"/api/images/{asset_id}")).status_code == 404
    finally:
        async with maker() as s:
            await s.execute(delete(User).where(User.email == email))
            await s.commit()


async def test_deleting_the_document_takes_the_files_too(api_client, maker, uploaded_docx):
    """⭐ 不能留孤儿：行、块、磁盘文件一起走。"""
    from sqlalchemy import select

    from copilot.db.models import ImageAsset

    async with maker() as s:
        paths = list(
            (
                await s.execute(
                    select(ImageAsset.storage_path).where(
                        ImageAsset.document_id == uploaded_docx
                    )
                )
            ).scalars()
        )
    assert paths

    r = await api_client.delete(f"/api/documents/{uploaded_docx}")
    assert r.status_code == 204, r.text

    async with maker() as s:
        left = list(
            (
                await s.execute(
                    select(ImageAsset.id).where(ImageAsset.document_id == uploaded_docx)
                )
            ).scalars()
        )
    assert left == []
    for rel in paths:
        assert not (get_settings().private_image_dir / rel).exists(), f"孤儿文件还在：{rel}"


async def test_a_shared_file_survives_another_document_being_deleted(
    api_client, maker, uploaded_docx, logged_in, flagship_id
):
    """⭐⭐ 同一张图被两篇文档引用时，删掉一篇**不能**把文件删掉。

    图是按内容寻址存的（一份文件、两行资产）。不数引用就删的话，
    删 A 会把 B 的图弄没——而 B 的表现是「图裂了」，没有任何报错，
    也再也恢复不回来（原始文件在用户电脑上）。
    """
    from sqlalchemy import delete, select

    from copilot.db.models import Document, ImageAsset

    async with maker() as s:
        rel = (
            await s.execute(
                select(ImageAsset.storage_path)
                .where(ImageAsset.document_id == uploaded_docx)
                .limit(1)
            )
        ).scalar_one()
        # 另一篇文档也引用同一个文件
        other = Document(
            owner_id=logged_in,
            knowledge_space_id=flagship_id,
            source_type="upload",
            title=f"引用同一张图-{uuid.uuid4().hex[:8]}",
            content_hash=uuid.uuid4().hex,
            status="done",
        )
        s.add(other)
        await s.flush()
        s.add(
            ImageAsset(
                document_id=other.id,
                owner_id=logged_in,
                knowledge_space_id=flagship_id,
                storage_path=rel,
                mime_type="image/png",
            )
        )
        await s.commit()
        other_id = other.id

    try:
        assert (await api_client.delete(f"/api/documents/{uploaded_docx}")).status_code == 204
        assert (get_settings().private_image_dir / rel).is_file(), "把还有人引用的图删掉了"
    finally:
        async with maker() as s:
            await s.execute(delete(ImageAsset).where(ImageAsset.document_id == other_id))
            await s.execute(delete(Document).where(Document.id == other_id))
            await s.commit()
