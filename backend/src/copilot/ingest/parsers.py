"""把用户上传的文件解析成 Markdown，交给既有的切分管线。

**为什么统一转 Markdown 而不是各格式各写一套切分**：`ingest/chunker.py` 的
分段能力全靠 Markdown 标题（`#`）——章节路径就是引用里那句
「订单对接-拼多多 · 第 2 节 授权与下载」的来源。所以解析器的首要职责不是
「把字取出来」，而是**把原文的层级还原成 `#` 标题**。docx 的
「标题 1」样式丢掉的话，一篇 30 页的操作手册会变成一坨没有章节的长文本，
检索命中了也说不出答案在哪一节。

服务器只有 1.6GB 内存，所以这里**只用纯 Python 的轻量库**
（python-docx / python-pptx / pypdf，都在 `parse` 这个 extra 里）。
PDF 只做纯文本提取，不上 Docling 那条会拖进 torch 的 ML 管线——
见 plan.md「一、第 3 条硬约束」。

**文档里的嵌图（M17）**：DOCX / PPTX / XLSX / PDF 里的截图要跟着正文一起
出来，而且**要落在它说明的那一段旁边**——ERP 操作手册的图就是步骤本身，
统一堆到文末的话，「第二步长什么样」这个问题就答不了了。

解析器自己不落盘、也不知道图片存到哪里：它把字节交给一个 `ImageSink`
（见 `copilot.assets.UploadImageSink`），拿回一个能写进 Markdown 的地址。
这样 parsers 保持可离线单测，而「私有图必须落进私有目录」这条规则
只有一处实现。`sink=None` 时所有格式**照旧只出文字**，一行都不改。

**图片和扫描件走另一条路**：本地跑不了 OCR（同样是 1.6GB 的约束），
所以送给视觉模型（Kimi）读成文字。这条路和上面几个解析器有一个本质区别——
**它要联网、要花钱、还可能编造内容**。所以：

- 视觉客户端是**注入**进来的（`parse_upload(..., vision=...)`），不在这里 import。
  parsers 保持可离线单测，也不会让「解析一个 txt」意外拖起一个 HTTP 客户端。
- 扫描件 PDF 有页数上限（`vision_pdf_max_pages`），它是一道**花钱的闸门**。
"""

from __future__ import annotations

import io
import re
from concurrent.futures import ThreadPoolExecutor
from concurrent.futures import TimeoutError as FutureTimeout
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol

from copilot.ingest.zipguard import UnsafeArchive, check_zip_safety

# docx 的样式名在中文版 Word 里是「标题 1」，英文版是 "Heading 1"。
# 两种都认，否则中文用户传上来的文档一个标题都识别不出。
_HEADING_STYLE_RE = re.compile(r"^(?:Heading|标题)\s*([1-6])$", re.I)
# Word 大纲级别也可能落在 "Title"/"标题"（无数字）上，那是文档大标题
_TITLE_STYLES = {"title", "标题", "文档标题"}


class ImageSink(Protocol):
    """解析器往里丢图片的口子。实现在 `copilot.assets.UploadImageSink`。

    `save` 返回**正文里该写的地址**；被闸门挡掉（太大、太小、超出张数上限）
    时返回 None——这时解析器**不要往正文里写这张图**：写一个取不到的地址
    比没有图更糟，页面上是一张裂图，而用户不知道为什么。
    """

    def save(self, data: bytes, suffix: str, **meta) -> str | None: ...


def _image_md(ref: str | None, alt: str = "") -> str:
    """图片的 Markdown。`ref` 为 None（被闸门挡掉）时什么都不写。"""
    return f"![{alt}]({ref})" if ref else ""


class ParseError(Exception):
    """解析失败。消息会原样存进 `documents.error` 给用户看，所以要写成人话。"""


@dataclass(slots=True)
class ParsedUpload:
    markdown: str
    #  给用户看的说明，比如「PDF 共 12 页」。存进 error 字段没意义，只用于日志
    note: str = ""


# ---------- 纯文本 ----------


def _decode(raw: bytes) -> str:
    """把字节解成文本。

    ⚠️ **必须有 GBK 兜底。** 中文 Windows 上记事本存的 .txt 默认不是 UTF-8，
    直接 `decode("utf-8")` 会抛 UnicodeDecodeError；而用 `errors="replace"`
    则更糟——不报错，但整篇变成一串「锟斤拷」，照样入库、照样检索不到，
    没有任何地方会提示出错。
    """
    for enc in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    raise ParseError("文件编码无法识别，请另存为 UTF-8 后重试")


def parse_text(path: Path) -> ParsedUpload:
    return ParsedUpload(markdown=_decode(path.read_bytes()))


# ---------- docx ----------


def _docx_table_to_markdown(table) -> str:
    """Word 表格转 Markdown 表格。

    ERP 的操作手册里大量内容在表格里（字段说明、对接映射表）。表格拍平成
    空格分隔的文字后，「哪个字段对应哪个含义」就丢了；Markdown 表格能保住
    这层对应关系，且模型读得懂。
    """
    rows: list[list[str]] = []
    for row in table.rows:
        # 单元格里的换行会把 Markdown 表格撑破，压成空格
        cells = [" ".join(c.text.split()) or "　" for c in row.cells]
        if any(c.strip() and c != "　" for c in cells):
            rows.append(cells)
    if not rows:
        return ""

    width = max(len(r) for r in rows)
    rows = [r + ["　"] * (width - len(r)) for r in rows]
    head, *body = rows
    lines = ["| " + " | ".join(head) + " |", "|" + "---|" * width]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)


# DOCX 里的图片：`a:blip` 上的 `r:embed` 指向关系 id，关系表再指到
# `word/media/*` 的实际字节。命名空间写全，别用 local-name() 那种模糊匹配
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _docx_paragraph_images(para, document, sink: ImageSink | None) -> list[str]:
    """这一段里嵌的图，按出现顺序。返回若干条 Markdown。

    ⚠️ **按段落取，不是一次性把 word/media 全捞出来。** 后者拿到的是一堆
    没有位置的图，只能统一贴到文末——而 ERP 手册里图就是步骤本身，
    「第二步长什么样」会因此答不了。
    """
    if sink is None:
        return []
    out: list[str] = []
    for blip in para._element.iter(f"{{{_A_NS}}}blip"):
        rid = blip.get(f"{{{_R_NS}}}embed")
        if not rid:
            continue
        try:
            part = document.part.related_parts[rid]
            data, ext = part.blob, Path(part.partname).suffix
        except (KeyError, AttributeError):
            # 关系表里找不到这张图（外链图片、或者文件本身不完整）。
            # 跳过一张图不该让整篇文档解析失败
            continue
        if md := _image_md(sink.save(data, ext or ".png")):
            out.append(md)
    return out


def parse_docx(path: Path, sink: ImageSink | None = None) -> ParsedUpload:
    """.docx → Markdown。标题样式转成 `#`，表格转成 Markdown 表格，嵌图跟着段落走。

    ⚠️ **必须按文档顺序遍历 body 的 XML 子节点。** python-docx 的
    `document.paragraphs` 和 `document.tables` 是两个独立列表，各自有序但
    彼此的相对位置丢了。用它们拼出来的结果是「全部段落，然后全部表格」——
    表格会整体跑到文末，脱离它所属的那一节。
    """
    try:
        import docx  # python-docx
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as e:  # pragma: no cover - 环境缺 extra 时才走到
        raise ParseError("服务端缺少 docx 解析组件（pip install '.[parse]'）") from e

    # ⚠️ **先过 ZIP 安全检查，再交给 python-docx**（M13 P9）。
    # .docx 本质是 ZIP，20MB 的上传上限管不住它解压出来有多大；
    # 而 worker 只有一个进程、`MemoryMax=400M`，一份炸弹就能把它拖死。
    # 顺序不能反：python-docx 一旦开始读，炸弹就已经在里面了
    try:
        check_zip_safety(path, kind="Word")
    except UnsafeArchive as e:
        raise ParseError(str(e)) from e

    try:
        document = docx.Document(str(path))
    except Exception as e:  # noqa: BLE001 - python-docx 抛的异常类型很杂
        raise ParseError(f"打不开这个 Word 文件（{type(e).__name__}）") from e

    parts: list[str] = []
    for child in document.element.body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "p":
            para = Paragraph(child, document)
            text = para.text.strip()
            style = (para.style.name or "").strip() if para.style is not None else ""
            images = _docx_paragraph_images(para, document, sink)
            if not text and not images:
                continue
            if m := _HEADING_STYLE_RE.match(style):
                parts.append(f"{'#' * int(m.group(1))} {text}")
            elif style.lower() in _TITLE_STYLES:
                parts.append(f"# {text}")
            elif text:
                parts.append(text)
            # 图紧跟在它所属的那一段后面
            parts.extend(images)
        elif tag == "tbl":
            if md := _docx_table_to_markdown(Table(child, document)):
                parts.append(md)

    markdown = "\n\n".join(parts).strip()
    if not markdown:
        raise ParseError("这个 Word 文件里没有可提取的文字（可能整篇都是图片）")
    return ParsedUpload(markdown=markdown)


# ---------- pptx ----------


def parse_pptx(path: Path, sink: ImageSink | None = None) -> ParsedUpload:
    """.pptx → Markdown，一页一节，图绑在它所在的那一页。

    每页做成一个 `##` 小节，标题优先用页面自己的标题占位符。PPT 的语义单元
    就是「页」，按页分节后引用能落到「第 7 页 · 面单打印设置」，比整篇一块有用得多。
    """
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise ParseError("服务端缺少 pptx 解析组件（pip install '.[parse]'）") from e

    # 同 parse_docx：.pptx 也是 ZIP，先过安全检查（M13 P9）
    try:
        check_zip_safety(path, kind="PPT")
    except UnsafeArchive as e:
        raise ParseError(str(e)) from e

    try:
        prs = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"打不开这个 PPT 文件（{type(e).__name__}）") from e

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = " ".join(slide.shapes.title.text.split())
        except (AttributeError, ValueError):
            title = ""  # 没有标题占位符的版式

        body: list[str] = []
        for shape in slide.shapes:
            # ⭐ 图片绑到**这一页**的小节里（`slide_number`）。绑错页的表现是
            # 答案配了另一页的截图——用户照着点会点不到，而没有任何报错
            image = getattr(shape, "image", None)
            if image is not None and sink is not None:
                try:
                    data, ext = image.blob, f".{(image.ext or 'png').lstrip('.')}"
                except (AttributeError, ValueError):
                    data = None
                if data and (md := _image_md(sink.save(data, ext, slide_number=i))):
                    body.append(md)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text or " ".join(text.split()) == title:
                continue
            body.append(text)

        if not title and not body:
            continue
        head = f"## 第 {i} 页 {title}".rstrip()
        parts.append(f"{head}\n\n" + "\n\n".join(body) if body else head)

    markdown = "\n\n".join(parts).strip()
    if not markdown:
        raise ParseError("这个 PPT 里没有可提取的文字（可能整篇都是图片）")
    return ParsedUpload(markdown=markdown, note=f"共 {len(prs.slides)} 页")


# ---------- 图片（视觉模型）----------


def parse_image(path: Path, vision=None) -> ParsedUpload:
    """.png/.jpg/... → Markdown，靠视觉模型读。

    没配 `VISION_API_KEY` 时报一句人话就停下。**不能默默入库一篇空文档**：
    那样用户看到的是「已完成」但永远搜不到，和上传坏了完全无法区分。
    """
    if vision is None:
        raise ParseError("服务端没有配置图片解析（VISION_API_KEY），暂时无法处理图片")

    from copilot.providers.vision import VisionError, mime_for

    try:
        text = vision.transcribe(path.read_bytes(), mime_for(path), hint=path.stem)
    except VisionError as e:
        raise ParseError(str(e)) from e

    if not text.strip():
        # 模型明确说了「这张图没有文字」。这不是失败，但也没有入库的价值——
        # 一篇没有内容的文档只会占着列表、混进检索
        raise ParseError("这张图片里没有可识别的文字内容")

    # 图片本身没有标题层级，补一个 `#`：切分器靠它给出「第 N 节」，
    # 引用里才不会是一句光秃秃的文件名
    title = path.stem
    if not text.lstrip().startswith("#"):
        text = f"# {title}\n\n{text}"
    return ParsedUpload(markdown=text, note="图片转写")


def _pdf_page_images(path: Path, pages: list[int], dpi: int) -> list[bytes]:
    """把指定页渲染成 PNG 字节。扫描件 PDF 专用。

    ⭐ 用 pypdfium2 不用 PyMuPDF：后者是 AGPL（plan.md 七.3 的许可红线）。

    **一页一渲染、渲完立刻关掉。** 一次性渲 20 页 A4@150dpi 是 130MB 的
    位图常驻，而 worker 的 `MemoryMax=400M`——那会让解析在中途被 systemd
    杀掉，用户看到的是「解析失败」而没有任何原因。
    """
    try:
        import pypdfium2 as pdfium
    except ImportError as e:  # pragma: no cover
        raise ParseError("服务端缺少 PDF 渲染组件（pip install '.[parse]'）") from e

    out: list[bytes] = []
    doc = pdfium.PdfDocument(str(path))
    try:
        for i in pages:
            page = doc[i]
            bitmap = page.render(scale=dpi / 72)
            img = bitmap.to_pil()
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            out.append(buf.getvalue())
            img.close()
            bitmap.close()
            page.close()
    finally:
        doc.close()
    return out


# ---------- xlsx ----------


# 一个工作表最多读多少行。十万行的流水表转成 Markdown 有几十 MB，
# 既撑爆 worker 的 400M 内存，也不可能是「知识」
XLSX_MAX_ROWS = 2000


def parse_xlsx(path: Path, sink: ImageSink | None = None) -> ParsedUpload:
    """.xlsx → Markdown，一个工作表一节，嵌图绑到它所在的工作表。

    ⚠️ **只读值，不读公式**（`data_only=True`）：库里要存的是「这个字段填
    什么」，不是「它是怎么算出来的」。读公式的话，一张对照表会变成一堆
    `=VLOOKUP(...)`，检索出来对谁都没用。

    ⚠️ **嵌图走 openpyxl 的 `ws._images`（私有属性），所以这一档标注为
    「有限支持」**（路线图第 37 节）：它不是公开 API，换个版本可能就没了。
    但为它引一个重量级依赖不划算——真没了就退回只出文字，
    下面这个 try 就是干这个的，不会让整篇解析失败。

    行数上限是刻意的：一张十万行的流水表转成 Markdown 有几十 MB，
    既撑爆 worker 的内存，也不可能是「知识」。
    """
    try:
        from openpyxl import load_workbook
    except ImportError as e:  # pragma: no cover
        raise ParseError("服务端缺少 Excel 解析组件（pip install '.[agent]'）") from e

    # 同 docx/pptx：.xlsx 也是 ZIP，先过安全检查（M13 P9）
    try:
        check_zip_safety(path, kind="Excel")
    except UnsafeArchive as e:
        raise ParseError(str(e)) from e

    try:
        wb = load_workbook(str(path), data_only=True, read_only=False)
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"打不开这个 Excel 文件（{type(e).__name__}）") from e

    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            rows: list[list[str]] = []
            for row in ws.iter_rows(max_row=XLSX_MAX_ROWS, values_only=True):
                cells = ["" if v is None else " ".join(str(v).split()) for v in row]
                if any(cells):
                    rows.append(cells)

            section = [f"## {ws.title}"]
            if rows:
                width = max(len(r) for r in rows)
                rows = [r + [""] * (width - len(r)) for r in rows]
                head, *body = rows
                table = ["| " + " | ".join(c or "　" for c in head) + " |", "|" + "---|" * width]
                table += ["| " + " | ".join(c or "　" for c in r) + " |" for r in body]
                section.append("\n".join(table))

            # 图片：`anchor` 记的是它钉在哪个单元格（"B7"），排查时有用
            for img in list(getattr(ws, "_images", []) or []):
                if sink is None:
                    break
                try:
                    data = img._data()
                    anchor = getattr(getattr(img, "anchor", None), "_from", None)
                    cell = (
                        f"R{anchor.row + 1}C{anchor.col + 1}" if anchor is not None else None
                    )
                except Exception:  # noqa: BLE001 - 私有 API，坏了就当没有图
                    continue
                if md := _image_md(
                    sink.save(data, ".png", sheet_name=ws.title[:128], anchor=cell)
                ):
                    section.append(md)

            if len(section) > 1:
                parts.append("\n\n".join(section))
    finally:
        wb.close()

    markdown = "\n\n".join(parts).strip()
    if not markdown:
        raise ParseError("这个 Excel 里没有可提取的内容")
    return ParsedUpload(markdown=markdown, note=f"共 {len(wb.worksheets)} 个工作表")


# ---------- pdf ----------


def _pdf_embedded_images(page, index: int, sink: ImageSink | None) -> list[str]:
    """这一页里嵌的图。pypdf 已经把它们解成可直接落盘的字节。

    ⚠️ **只取嵌图，不做整页截图。** 路线图提过「每页生成 page snapshot」，
    但那要把每一页渲染成位图——一份 200 页的 PDF 就是 200 次渲染、几百 MB
    的中间位图，而 worker 的 MemoryMax 是 400M。数字版 PDF 的文字本来就
    提得出来，扫描件那条路（整页渲染 + 视觉识别）已经在下面了。
    """
    if sink is None:
        return []
    out: list[str] = []
    try:
        images = list(page.images)
    except Exception:  # noqa: BLE001 - 单页的图坏了不该毁掉整篇
        return []
    for img in images:
        try:
            data, name = img.data, (img.name or "")
        except Exception:  # noqa: BLE001
            continue
        if md := _image_md(sink.save(data, Path(name).suffix or ".png", page_number=index)):
            out.append(md)
    return out


def parse_pdf(path: Path, vision=None, sink: ImageSink | None = None) -> ParsedUpload:
    """.pdf → 纯文本，一页一节，嵌图绑到它所在的那一页。

    **只做纯文本提取**，不做版面还原：PDF 里没有「标题样式」这种结构信息，
    要还原层级得上 ML 版面分析（Docling），而那会拖进 torch——1.6GB 的机器
    装都装不下。所以 PDF 的检索质量天生弱于 docx，这是明确的取舍。

    ⚠️ **提不出字要报错，不能静默入库。** 扫描件 PDF（整页都是图）
    `extract_text()` 返回空串，不拦的话会入库一篇空文档，
    用户看到状态「已完成」却永远搜不到内容，也不知道为什么。
    """
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover
        raise ParseError("服务端缺少 PDF 解析组件（pip install '.[parse]'）") from e

    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            # 空密码能开的加密 PDF 不少（只设了权限密码），试一下再放弃
            try:
                reader.decrypt("")
            except Exception as e:  # noqa: BLE001
                raise ParseError("这个 PDF 有密码保护，请先解除后上传") from e
        pages = reader.pages
    except ParseError:
        raise
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"打不开这个 PDF 文件（{type(e).__name__}）") from e

    parts: list[str] = []
    for i, page in enumerate(pages, 1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:  # noqa: BLE001 - 单页坏了不该毁掉整篇
            text = ""
        images = _pdf_embedded_images(page, i, sink)
        if not text and not images:
            continue
        body = "\n\n".join([p for p in [text, *images] if p])
        parts.append(f"## 第 {i} 页\n\n{body}")

    markdown = "\n\n".join(parts).strip()
    if markdown:
        return ParsedUpload(markdown=markdown, note=f"共 {len(pages)} 页")

    # 一个字都提不出 = 扫描件（整页都是图）。交给视觉模型逐页读。
    #
    # ⚠️ **判据是「整篇都空」，不是「某几页空」。** 按页回退看着更聪明，
    # 实际上一份 200 页的正常 PDF 里夹几页插图是常态，那样会在用户毫不知情的
    # 情况下反复触发付费调用。宁可漏掉那几页插图。
    if vision is None:
        raise ParseError("这个 PDF 提取不出文字，可能是扫描件（图片型 PDF），暂不支持")
    return _parse_scanned_pdf(path, len(pages), vision)


def _parse_scanned_pdf(path: Path, total: int, vision) -> ParsedUpload:
    """扫描件 PDF：逐页渲染成图，再交给视觉模型。"""
    from copilot.config import get_settings
    from copilot.providers.vision import VisionError

    s = get_settings()
    limit = s.vision_pdf_max_pages
    todo = list(range(min(total, limit)))
    images = _pdf_page_images(path, todo, s.vision_pdf_dpi)

    parts: list[str] = []
    for i, raw in enumerate(images, 1):
        try:
            text = vision.transcribe(raw, "image/png", hint=f"{path.stem} 第 {i} 页")
        except VisionError as e:
            # 单页失败不该毁掉整篇——但**要留痕**。静默跳过的话，
            # 用户拿到的是一份缺了几页却看起来完整的文档
            parts.append(f"## 第 {i} 页\n\n[这一页没能识别：{e}]")
            continue
        if text.strip():
            parts.append(f"## 第 {i} 页\n\n{text}")

    markdown = "\n\n".join(parts).strip()
    if not markdown:
        raise ParseError("这份扫描件每一页都读不出文字，请确认图像是否清晰")

    note = f"扫描件，视觉识别 {len(todo)} 页"
    if total > limit:
        # 截断必须说出来。不说的话，用户搜不到第 21 页的内容时
        # 只会以为知识库不好用
        note += f"（共 {total} 页，超出上限 {limit} 页的部分未处理）"
    return ParsedUpload(markdown=markdown, note=note)


# ---------- 分派 ----------

PARSERS = {
    ".md": parse_text,
    ".txt": parse_text,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".pdf": parse_pdf,
}

# 认 `sink=` 这个参数的解析器。纯文本那两个不认——给它们传等于要求
# 每个解析器都长一样，而 .txt 里根本不可能有嵌图
_SINK_PARSERS = {parse_docx, parse_pptx, parse_xlsx}

# 这几个必须有视觉模型才解析得了。单列一份是为了让 `parse_upload` 能在
# 真正读文件之前就判断「这台机器行不行」
VISION_PARSERS = {
    ".png": parse_image,
    ".jpg": parse_image,
    ".jpeg": parse_image,
    ".webp": parse_image,
    ".bmp": parse_image,
}


# 本机解析的时间上限（秒）。M13 P9。
#
# ⭐ **它要挡的不是慢，是「永远不结束」。** worker 只有一个进程、解析是同步的
# CPU 活（故意没丢线程池，见 jobs/worker.py 文件头）——一份损坏或者刁钻构造的
# 文件让 python-docx 转不出来，后面所有人的上传就一直停在「解析中」，
# 而页面上没有任何异常提示。
#
# 120 秒是量出来的余量：本机实测一份 20MB（上传上限）的 pptx 走完
# 解析 + 转 Markdown 在个位数秒。留一个数量级以上，正常文档碰不到它。
PARSER_TIMEOUT = 120.0

# 超时时给用户看的话。**不要把 Python 的堆栈或者 `TimeoutError` 露出去**：
# 这句会原样存进 `documents.error`，显示在他的知识库页面上
_TIMEOUT_MESSAGE = "文件解析时间过长，可能文件损坏或内容异常。请检查文件后重新上传。"


def _run_with_timeout(fn, *args, timeout: float | None = None, **kwargs) -> ParsedUpload:
    """跑一个同步解析器，超时就放弃。

    ⚠️ **诚实地说清楚这一道的边界：Python 杀不掉一个正在跑的线程。**
    超时之后这里立刻返回（任务被判失败、队列继续往下走、用户看到一句人话），
    但那条线程**还在后台烧 CPU**，直到它自己跑完或者进程被换掉。

    所以这一道的作用是**把「无限期卡住」降级成「一次失败」**，不是真的回收
    资源。真正兜住资源的是另外两道，它们都已经在位：
        `zipguard` —— 绝大多数病态输入在解析开始之前就被挡掉了
        systemd `MemoryMax=400M` + `Restart=always` —— 真吃爆了被收走的是
        worker，网站还在，而且它会自己起来

    `shutdown(wait=False)`：不等那条线程，否则这里就白超时了。
    """
    # ⚠️ **在调用时才读 `PARSER_TIMEOUT`，不能把它写成默认参数值。**
    # 默认参数在 import 时求值一次就定死了——之后无论是测试 monkeypatch
    # 还是以后改成从配置读，改的都是一个再也没人看的模块变量，
    # 而超时依旧按 import 那一刻的数字走。这种失效是完全静默的。
    timeout = PARSER_TIMEOUT if timeout is None else timeout
    pool = ThreadPoolExecutor(max_workers=1, thread_name_prefix="parse")
    try:
        return pool.submit(fn, *args, **kwargs).result(timeout=timeout)
    except FutureTimeout as e:
        raise ParseError(_TIMEOUT_MESSAGE) from e
    finally:
        pool.shutdown(wait=False)


def parse_upload(
    path: Path, suffix: str | None = None, vision=None, sink: ImageSink | None = None
) -> ParsedUpload:
    """按扩展名解析。

    `suffix` 可显式指定——落盘名是 uuid，但后缀保留，一般不用传。
    `vision` 是读图客户端，只有图片和扫描件 PDF 用得上；传 None 时
    图片会得到一句「服务端没有配置图片解析」，而不是一个 AttributeError。

    ⚠️ **超时只套在本机解析那几个上，不套视觉那条路**（M13 P9）。
    视觉路是网络调用：一份 20 页的扫描件要逐页发给 Kimi，一页几秒，
    一份合法文件就能超过 120 秒。拿本机 CPU 的口径去卡它，
    只会把「正常但慢」判成「文件损坏」。那条路自己有 httpx 的超时，
    还有 `vision_pdf_max_pages` 这道花钱的闸门。
    """
    ext = (suffix or path.suffix).lower()
    if ext in VISION_PARSERS:
        return VISION_PARSERS[ext](path, vision=vision)
    parser = PARSERS.get(ext)
    if parser is None:
        raise ParseError(f"不支持的文件类型：{ext or '（无扩展名）'}")
    if parser is parse_pdf:
        # 扫描件 PDF 会走进视觉那条路，同样不能卡本机口径
        return parse_pdf(path, vision=vision, sink=sink)
    if parser in _SINK_PARSERS:
        return _run_with_timeout(parser, path, sink=sink)
    return _run_with_timeout(parser, path)
